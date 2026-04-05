import json
import re
from typing import Any

from crewai import Agent, Crew, LLM, Task
from pptx import Presentation

from app.config import get_settings

settings = get_settings()


def extract_ppt_text(ppt_path: str) -> str:
    """Extract all readable text from a .pptx file."""
    if not ppt_path.lower().endswith(".pptx"):
        raise ValueError("Only .pptx files are supported")

    presentation = Presentation(ppt_path)
    slides_text: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)

        if slide_parts:
            slides_text.append(f"[Slide {index}]\n" + "\n".join(slide_parts))

    return "\n\n".join(slides_text)


def _extract_json_from_text(text: str) -> dict[str, Any]:
    """Extract JSON safely from model output."""
    cleaned = text.strip()

    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    fenced_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", cleaned, re.DOTALL)
    if fenced_match:
        return json.loads(fenced_match.group(1))

    json_match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if json_match:
        return json.loads(json_match.group(0))

    raise ValueError("No valid JSON object found in model response")


def _normalize_scores(
    result: dict[str, Any], rubric_criteria: list[dict[str, Any]]
) -> dict[str, Any]:
    """Clamp scores to rubric max_score and keep output consistent."""
    max_score_map = {
        item["category"]: float(item["max_score"])
        for item in rubric_criteria
        if "category" in item and "max_score" in item
    }

    normalized_scores: list[dict[str, Any]] = []

    for rubric_item in rubric_criteria:
        category = rubric_item["category"]
        max_score = float(rubric_item["max_score"])

        matched = None
        for score_item in result.get("criteria_scores", []):
            if score_item.get("category") == category:
                matched = score_item
                break

        if matched is None:
            normalized_scores.append(
                {
                    "category": category,
                    "score": 0.0,
                    "comment": "No evaluation returned for this criterion.",
                }
            )
            continue

        raw_score = matched.get("score", 0)
        try:
            numeric_score = float(raw_score)
        except (TypeError, ValueError):
            numeric_score = 0.0

        numeric_score = max(0.0, min(numeric_score, max_score))

        normalized_scores.append(
            {
                "category": category,
                "score": numeric_score,
                "comment": str(matched.get("comment", "")).strip(),
            }
        )

    result["criteria_scores"] = normalized_scores
    result["ppt_summary"] = str(result.get("ppt_summary", "")).strip()

    return result


def analyze_ppt(ppt_path: str, rubric_criteria: list[dict[str, Any]]) -> dict[str, Any]:
    """
    Analyze PPT content against rubric criteria using CrewAI + Gemini.
    Expects a local .pptx path.
    """
    try:
        slide_text = extract_ppt_text(ppt_path)

        if not slide_text.strip():
            return {
                "criteria_scores": [],
                "ppt_summary": "No readable text was found in the presentation.",
                "error": "Empty presentation content",
            }

        criteria_text = "\n".join(
            [
                f"- {item['category']} (max {item['max_score']} pts): {item['description']}"
                for item in rubric_criteria
            ]
        )

        llm = LLM(
                model="gemini-2.5-flash",
                api_key=settings.GEMINI_API_KEY,
            )

        ppt_analyzer_agent = Agent(
            role="PPT Analyzer",
            goal="Evaluate PowerPoint slides against rubric criteria and return strict JSON output.",
            backstory=(
                "You are an academic presentation evaluator. "
                "You score each rubric criterion carefully, justify each score briefly, "
                "and always return valid JSON only."
            ),
            llm=llm,
            verbose=False,
            allow_delegation=False,
        )

        prompt = f"""
You are evaluating an academic presentation.

RUBRIC CRITERIA:
{criteria_text}

PRESENTATION CONTENT:
{slide_text}

Instructions:
- Evaluate each rubric criterion separately.
- Score each criterion from 0 up to that criterion's max_score.
- Give a short but clear explanation for each criterion.
- Give an overall 2-3 sentence summary of the presentation.
- Return ONLY valid JSON.
- Do not include markdown.
- Do not include any extra text before or after the JSON.

Return exactly this JSON structure:
{{
  "criteria_scores": [
    {{
      "category": "<criterion name>",
      "score": <number>,
      "comment": "<brief explanation>"
    }}
  ],
  "ppt_summary": "<overall 2-3 sentence summary>"
}}
"""

        analyze_task = Task(
            description=prompt,
            expected_output="A valid JSON object with criterion-wise scores and a presentation summary.",
            agent=ppt_analyzer_agent,
        )

        crew = Crew(
            agents=[ppt_analyzer_agent],
            tasks=[analyze_task],
            verbose=False,
        )

        result = crew.kickoff()
        parsed = _extract_json_from_text(str(result))
        normalized = _normalize_scores(parsed, rubric_criteria)
        return normalized

    except Exception as exc:
        return {
            "criteria_scores": [],
            "ppt_summary": "",
            "error": f"PPT analysis failed: {str(exc)}",
        }
