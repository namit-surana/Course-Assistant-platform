from __future__ import annotations

from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
from pptx import Presentation


def _default_rubric_path() -> Path:
    return Path(__file__).resolve().parent / "knowledge" / "default_rubric.md"


def load_default_ppt_rubric() -> str:
    """Course-standard presentation evaluation guidance (markdown), like the video agent rubric."""
    return _default_rubric_path().read_text(encoding="utf-8")


# Fixed scoring dimensions (aligned with knowledge/default_rubric.md). Worker does not use professor rubrics.
BUILTIN_PPT_RUBRIC_CRITERIA: list[dict[str, Any]] = [
    {"category": "A1", "max_score": 5.0, "description": "Document is readable: substantive text, not blank or placeholder-only."},
    {"category": "A2", "max_score": 5.0, "description": "Project and team identifiable early if expected (titles, headers, first slides)."},
    {"category": "A3", "max_score": 5.0, "description": "Deck/PDF format appropriate; content maps to slides/pages consistently."},
    {"category": "B1", "max_score": 5.0, "description": "Logical flow: context → approach → results/demo → conclusion or next steps."},
    {"category": "B2", "max_score": 5.0, "description": "Sections or slide titles support the story; ordering makes sense."},
    {"category": "B3", "max_score": 5.0, "description": "Enough depth to follow the argument, not only buzzwords."},
    {"category": "C1", "max_score": 5.0, "description": "Features, architecture, or methodology described concretely."},
    {"category": "C2", "max_score": 5.0, "description": "Integrations/external systems explained or gaps acknowledged when relevant."},
    {"category": "C3", "max_score": 5.0, "description": "Limitations, mocks, or unfinished areas disclosed when relevant."},
    {"category": "D1", "max_score": 5.0, "description": "Terminology consistent; acronyms explained when needed."},
    {"category": "D2", "max_score": 5.0, "description": "Diagrams/lists referenced in text enough to understand intent without seeing the image."},
    {"category": "D3", "max_score": 5.0, "description": "Takeaways or conclusions explicit where expected."},
    {"category": "E1", "max_score": 5.0, "description": "Claims match detail shown; no unsupported overstated “done”."},
    {"category": "E2", "max_score": 5.0, "description": "Repo/demo/video links only if they appear in the extract."},
    {"category": "F1", "max_score": 5.0, "description": "Reproducibility or setup notes if applicable to the assignment type."},
    {"category": "F2", "max_score": 5.0, "description": "Timeline, risks, or teamwork only if present and relevant."},
]


def builtin_ppt_rubric_by_category() -> dict[str, dict[str, Any]]:
    return {str(row["category"]): row for row in BUILTIN_PPT_RUBRIC_CRITERIA}


def extract_ppt_text(file_path: str) -> str:
    """Extract readable text from a .pptx file."""
    presentation = Presentation(file_path)
    slides_text: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)

        if slide_parts:
            slides_text.append(f"[Slide {index}]\n" + "\n".join(slide_parts))

    return "\n\n".join(slides_text)


def extract_pdf_text(file_path: str) -> str:
    """Extract readable text from a .pdf file."""
    doc = fitz.open(file_path)
    pages_text: list[str] = []

    for index, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if text:
            pages_text.append(f"[Page {index}]\n{text}")

    doc.close()
    return "\n\n".join(pages_text)


def extract_document_text(file_path: str) -> str:
    """Extract text from supported document formats."""
    lower_path = file_path.lower()

    if lower_path.endswith(".pptx"):
        return extract_ppt_text(file_path)

    if lower_path.endswith(".pdf"):
        return extract_pdf_text(file_path)

    raise ValueError("Only .pptx and .pdf files are supported")

