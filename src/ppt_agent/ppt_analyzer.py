import json
import re
from pathlib import Path
from typing import Any

import fitz  # PyMuPDF
from crewai import Agent, Crew, LLM, Task
from pptx import Presentation

from src.config.settings import get_settings

settings = get_settings()


def _default_rubric_path() -> Path:
    return Path(__file__).resolve().parent / "knowledge" / "default_rubric.md"


def load_default_ppt_rubric() -> str:
    """Course-standard presentation evaluation guidance (markdown), like the video agent rubric."""
    return _default_rubric_path().read_text(encoding="utf-8")


# Fixed scoring dimensions (aligned with knowledge/default_rubric.md). Worker does not use professor rubrics.
BUILTIN_PPT_RUBRIC_CRITERIA: list[dict[str, Any]] = [
    {"category": "A1", "max_score": 5.0, "description": "Document is readable: substantive text, not blank or placeholder-only."},
    {"category": "A2", "max_score": 5.0, "description": "Project and team identifiable early if expected (titles, headers, first slides)."},
    {"category": "A3", "max_score": 5.0, "description": "Deck/PDF format appropriate; content maps to slides/pages consistently."},
    {"category": "B1", "max_score": 5.0, "description": "Logical flow: context → approach → results/demo → conclusion or next steps."},
    {"category": "B2", "max_score": 5.0, "description": "Sections or slide titles support the story; ordering makes sense."},
    {"category": "B3", "max_score": 5.0, "description": "Enough depth to follow the argument, not only buzzwords."},
    {"category": "C1", "max_score": 5.0, "description": "Features, architecture, or methodology described concretely."},
    {"category": "C2", "max_score": 5.0, "description": "Integrations/external systems explained or gaps acknowledged when relevant."},
    {"category": "C3", "max_score": 5.0, "description": "Limitations, mocks, or unfinished areas disclosed when relevant."},
    {"category": "D1", "max_score": 5.0, "description": "Terminology consistent; acronyms explained when needed."},
    {"category": "D2", "max_score": 5.0, "description": "Diagrams/lists referenced in text enough to understand intent without seeing the image."},
    {"category": "D3", "max_score": 5.0, "description": "Takeaways or conclusions explicit where expected."},
    {"category": "E1", "max_score": 5.0, "description": "Claims match detail shown; no unsupported overstated “done”."},
    {"category": "E2", "max_score": 5.0, "description": "Repo/demo/video links only if they appear in the extract."},
    {"category": "F1", "max_score": 5.0, "description": "Reproducibility or setup notes if applicable to the assignment type."},
    {"category": "F2", "max_score": 5.0, "description": "Timeline, risks, or teamwork only if present and relevant."},
]


def builtin_ppt_rubric_by_category() -> dict[str, dict[str, Any]]:
    return {str(row["category"]): row for row in BUILTIN_PPT_RUBRIC_CRITERIA}


def extract_ppt_text(file_path: str) -> str:
    """Extract readable text from a .pptx file."""
    presentation = Presentation(file_path)
    slides_text: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)

        if slide_parts:
            slides_text.append(f"[Slide {index}]\n" + "\n".join(slide_parts))

    return "\n\n".join(slides_text)


def extract_pdf_text(file_path: str) -> str:
    """Extract readable text from a .pdf file."""
    doc = fitz.open(file_path)
    pages_text: list[str] = []

    for index, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if text:
            pages_text.append(f"[Page {index}]\n{text}")

    doc.close()
    return "\n\n".join(pages_text)


def extract_document_text(file_path: str) -> str:
    """Extract text from supported document formats."""
    lower_path = file_path.lower()

    if lower_path.endswith(".pptx"):
        return extract_ppt_text(file_path)

    if lower_path.endswith(".pdf"):
        return extract_pdf_text(file_path)

    raise ValueError("Only .pptx and .pdf files are supported")


def _extract_json_from_text(text: str) -> dict[str, Any]:
    cleaned = text.strip()

    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        pass

    fenced_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", cleaned, re.DOTALL)
    if fenced_match:
        return json.loads(fenced_match.group(1))

    json_match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if json_match:
        return json.loads(json_match.group(0))

    raise ValueError("No valid JSON object found in model response")


def _normalize_scores(
    result: dict[str, Any], rubric_criteria: list[dict[str, Any]]
) -> dict[str, Any]:
    normalized_scores: list[dict[str, Any]] = []

    for rubric_item in rubric_criteria:
        category = rubric_item["category"]
        max_score = float(rubric_item["max_score"])

        matched = None
        for score_item in result.get("criteria_scores", []):
            if score_item.get("category") == category:
                matched = score_item
                break

        if matched is None:
            normalized_scores.append(
                {
                    "category": category,
                    "score": 0.0,
                    "comment": "No evaluation returned for this criterion.",
                }
            )
            continue

        try:
            numeric_score = float(matched.get("score", 0))
        except (TypeError, ValueError):
            numeric_score = 0.0

        numeric_score = max(0.0, min(numeric_score, max_score))

        normalized_scores.append(
            {
                "category": category,
                "score": numeric_score,
                "comment": str(matched.get("comment", "")).strip(),
            }
        )

    result["criteria_scores"] = normalized_scores
    result["ppt_summary"] = str(result.get("ppt_summary", "")).strip()

    return result


def analyze_ppt(
    file_path: str,
    rubric_criteria: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    """
    Analyze PPT/PDF content using CrewAI + Gemini.

    Pass a non-empty ``rubric_criteria`` only for ad-hoc tools (e.g. /analyze-ppt). The worker omits
    this so :data:`BUILTIN_PPT_RUBRIC_CRITERIA` is always used.
    """
    effective_criteria = rubric_criteria if rubric_criteria else BUILTIN_PPT_RUBRIC_CRITERIA
    try:
        document_text = extract_document_text(file_path)

        if not document_text.strip():
            return {
                "criteria_scores": [],
                "ppt_summary": "No readable text was found in the file.",
                "error": "Empty file content",
            }

        criteria_text = "\n".join(
            [
                f"- {item['category']} (max {item['max_score']} pts): {item['description']}"
                for item in effective_criteria
            ]
        )

        rubric_guidance = load_default_ppt_rubric()

        llm = LLM(
            model="gemini/gemini-2.5-flash-lite",
            api_key=settings.GEMINI_API_KEY,
        )

        analyzer_agent = Agent(
            role="Presentation Analyzer",
            goal="Evaluate presentation or document content against rubric criteria and return strict JSON output.",
            backstory=(
                "You are an academic evaluator. You score each rubric criterion carefully, "
                "justify each score briefly, and always return valid JSON only."
            ),
            llm=llm,
            verbose=False,
            allow_delegation=False,
        )

        prompt = f"""
You are evaluating an academic presentation/document.

EVALUATION GUIDANCE (how to read slides and write fair comments; align scores with this where it applies):
{rubric_guidance}

AUTHORITATIVE SCORING CRITERIA (you must score every row below; use category names exactly):
{criteria_text}

DOCUMENT CONTENT (extracted text only — you cannot see layout, animations, or images without OCR):
{document_text}

Instructions:
- Evaluate each line under AUTHORITATIVE SCORING CRITERIA separately; those names and max scores are binding.
- Use EVALUATION GUIDANCE to interpret thin extracts and to calibrate comments (do not invent slide content not in the extract).
- Score each criterion as an integer from 0 up to that criterion's max_score (no 0–1 normalization).
- Give a short but clear explanation for each criterion.
- Give an overall 2-3 sentence summary.
- Return ONLY valid JSON.
- Do not include markdown.
- Do not include any extra text before or after the JSON.

Return exactly this JSON structure:
{{
  "criteria_scores": [
    {{
      "category": "<criterion name>",
      "score": <number>,
      "comment": "<brief explanation>"
    }}
  ],
  "ppt_summary": "<overall 2-3 sentence summary>"
}}
"""

        analyze_task = Task(
            description=prompt,
            expected_output="A valid JSON object with criterion-wise scores and a summary.",
            agent=analyzer_agent,
        )

        crew = Crew(
            agents=[analyzer_agent],
            tasks=[analyze_task],
            verbose=False,
        )

        result = crew.kickoff()
        parsed = _extract_json_from_text(str(result))
        return _normalize_scores(parsed, effective_criteria)

    except Exception as exc:
        return {
            "criteria_scores": [],
            "ppt_summary": "",
            "error": f"Analysis failed: {str(exc)}",
        }
